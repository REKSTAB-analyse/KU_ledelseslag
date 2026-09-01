import io

import streamlit as st
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
import matplotlib
#matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from config import ROOT_ID, ROOT_NAVN, NIVEAUER, load_real_units
from data.loader import load_logo, logo_base64


PPTX_SKABELON = r"C:\Users\rjp530\Desktop\kontormøde_pptx\ku_skabelon.pptx"
PPTX_LAYOUT_NAVN = "1_Title and Content"

@st.cache_data
def load_units():
    return load_real_units()

def build_lookup_and_rollup(units):
    """
    by_id: {id: enhed-dict}, inkl. en tilføjet rod-enhed (ROOT_ID).
    children_of: {parent_id: [child_id, ...]}
    Aarsvaerk/lonomkostninger/medarbejdere rulles op, så ALLE enheder (ikke
    kun leaf-enheder) har summerede tal.
    """
    by_id = {u["id"]: dict(u) for u in units}
    by_id[ROOT_ID] = {
        "id": ROOT_ID, "navn": ROOT_NAVN, "niveau": "Rod",
        "parent_id": None, "ledelseslag": "Rektorat/direktion",
        "aarsvaerk": None, "lonomkostninger": None, "medarbejdere": None,
    }
 
    children_of = {}
    for u in by_id.values():
        if u["parent_id"] is not None:
            children_of.setdefault(u["parent_id"], []).append(u["id"])
 
    def rollup(unit_id):
        u = by_id[unit_id]
        kids = children_of.get(unit_id, [])
        if not kids:
            return u["aarsvaerk"] or 0.0, u["lonomkostninger"] or 0, u.get("medarbejdere") or 0
        total_av, total_lon, total_med = 0.0, 0, 0
        for k in kids:
            av, lon, med = rollup(k)
            total_av += av
            total_lon += lon
            total_med += med
        u["aarsvaerk"] = round(total_av, 1)
        u["lonomkostninger"] = total_lon
        u["medarbejdere"] = total_med
        return total_av, total_lon, total_med
 
    rollup(ROOT_ID)
    return by_id, children_of
 
 
def gns_loen(by_id, unit_id):
    u = by_id[unit_id]
    if not u["aarsvaerk"]:
        return 0
    return u["lonomkostninger"] / u["aarsvaerk"]
 
 
def path_to_root(by_id, unit_id):
    """Brødkrumme fra roden ned til unit_id, som liste af id'er."""
    path = [unit_id]
    while by_id[path[-1]]["parent_id"] is not None:
        path.append(by_id[path[-1]]["parent_id"])
    return list(reversed(path))
 
 
def leaves_under(children_of, unit_id):
    """
    Alle leaf-enheder (det yderste niveau, fx Kontor) under unit_id,
    fladtgjort på tværs af varierende dybde. Har unit_id selv ingen børn,
    returneres en tom liste.
    """
    kids = children_of.get(unit_id, [])
    if not kids:
        return []
    leaves = []
    for k in kids:
        if children_of.get(k):
            leaves.extend(leaves_under(children_of, k))
        else:
            leaves.append(k)
    return leaves

def _split_by_omraade(by_id, children_of, enh_uid, omraade_valgt, metric):
    """
    Deler en enheds kontorer i to grupper - dem der hører til omraade_valgt,
    og resten - og returnerer (omraade_vaerdi, rest_vaerdi) for den valgte
    metric. For de to "pr. X"-metrics beregnes et separat gennemsnit pr.
    gruppe (kan ikke stakkes, da et gennemsnit ikke er additivt) - for de
    øvrige (rene sum-metrics) summeres der (additive, kan stakkes).
    """
    kontor_ids = children_of.get(enh_uid, [])
    om_ids = [k for k in kontor_ids if by_id[k]["omraade"] == omraade_valgt]
    rest_ids = [k for k in kontor_ids if by_id[k]["omraade"] != omraade_valgt]

    if metric == "Gns. lønomkostning pr. årsværk":
        om_av = sum(by_id[k]["aarsvaerk"] for k in om_ids)
        rest_av = sum(by_id[k]["aarsvaerk"] for k in rest_ids)
        om_v = (sum(by_id[k]["lonomkostninger"] for k in om_ids) / om_av) if om_av else 0
        rest_v = (sum(by_id[k]["lonomkostninger"] for k in rest_ids) / rest_av) if rest_av else 0
    elif metric == "Gns. lønomkostning pr. medarbejder":
        om_med = sum(by_id[k]["medarbejdere"] for k in om_ids)
        rest_med = sum(by_id[k]["medarbejdere"] for k in rest_ids)
        om_v = (sum(by_id[k]["lonomkostninger"] for k in om_ids) / om_med) if om_med else 0
        rest_v = (sum(by_id[k]["lonomkostninger"] for k in rest_ids) / rest_med) if rest_med else 0
    elif metric == "Samlede lønomkostninger":
        om_v = sum(by_id[k]["lonomkostninger"] for k in om_ids)
        rest_v = sum(by_id[k]["lonomkostninger"] for k in rest_ids)
    elif metric == "Antal medarbejdere":
        om_v = sum(by_id[k]["medarbejdere"] for k in om_ids)
        rest_v = sum(by_id[k]["medarbejdere"] for k in rest_ids)
    else:  # "Antal årsværk"
        om_v = sum(by_id[k]["aarsvaerk"] for k in om_ids)
        rest_v = sum(by_id[k]["aarsvaerk"] for k in rest_ids)
    return om_v, rest_v

def _bar_chart_png(navne, values, farve_hex, metric_label, vaerdi_er_kr, width_in, height_in):
    """Bygger et vandret søjlediagram med matplotlib - samme stil som appens
    egne Plotly-diagrammer (KU-farver, størst øverst) - og returnerer det
    som PNG-bytes i en BytesIO. Kræver ikke Chrome."""
    farver = [farve_hex] * len(navne) if isinstance(farve_hex, str) else farve_hex
    colors = [f"#{f}" for f in farver]
    
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=200)
    y_pos = range(len(navne))
    ax.barh(y_pos, values, color=colors)
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(navne, fontsize=9)
    ax.invert_yaxis()
    ax.set_xlabel(metric_label, fontsize=9)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    if vaerdi_er_kr:
        ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}".replace(",", ".")))

    for i, v in enumerate(values):
        label = f"{v:,.0f} kr.".replace(",", ".") if vaerdi_er_kr else f"{v:.1f} årsværk"
        ax.text(v, i, " " + label, va="center", fontsize=8)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png")
    plt.close(fig)
    buf.seek(0)
    return buf

def build_full_pptx(by_id, children_of, niveau1_ids, metric, metric_value):
    """
    Bygger en pptx ud fra KU-skabelonen med native PowerPoint-diagrammer
    (kræver ikke Chrome/kaleido, i modsætning til billedeksport af plotly-
    figurer). Slide 1: oversigt over alle CA/KE-enheder. Slide 2-N: én
    slide PR. ENHED med to diagrammer side om side - overblikket til
    venstre, enhedens kontorer til højre - dvs. alle enheder, ikke kun
    den der aktuelt er valgt/vist i appen.
    """
    prs = Presentation(PPTX_SKABELON)
    layout = next((l for l in prs.slide_layouts if l.name == PPTX_LAYOUT_NAVN), None)
    if layout is None:
        layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]

    vaerdi_er_kr = metric != "Antal årsværk"

    def set_title(slide, title):
        title_ph = next((p for p in slide.placeholders if p.placeholder_format.idx == 0), None)
        if title_ph is not None:
            title_ph.text_frame.text = title
        else:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
            box.text_frame.text = title

    # Slide 1: alle CA/KE-enheder
    navne = [by_id[uid]["navn"] for uid in niveau1_ids]
    values = [metric_value(uid) for uid in niveau1_ids]

    slide = prs.slides.add_slide(layout)
    set_title(slide, "Campusadministrationer og koncernenheder")
    def add_chart_image(slide, left, top, width, height, navne, values, farve_hex):
        png_buf = _bar_chart_png(
            navne, values, farve_hex, metric, vaerdi_er_kr,
            width_in=width.inches, height_in=height.inches,
        )
        slide.shapes.add_picture(png_buf, left, top, width=width, height=height)
    add_chart_image(
        slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.4),
        navne, values, "901A1E",
    )

    # Slide 2..N: én slide pr. enhed, med to diagrammer side om side
    for uid in niveau1_ids:
        leaf_ids = leaves_under(children_of, uid)
        if not leaf_ids:
            continue
        leaf_ids = sorted(leaf_ids, key=lambda lid: by_id[lid]["navn"])
        leaf_navne = [by_id[lid]["navn"] for lid in leaf_ids]
        leaf_values = [metric_value(lid) for lid in leaf_ids]

        slide = prs.slides.add_slide(layout)
        set_title(slide, by_id[uid]["navn"])
        bar_farver = ["901A1E" if uid2 == uid else "E6C9CC" for uid2 in niveau1_ids]
        add_chart_image(
            slide, Inches(0.4), Inches(1.7), Inches(6.1), Inches(5.4),
            navne, values, bar_farver,
        )
        add_chart_image(
            slide, Inches(6.8), Inches(1.7), Inches(6.1), Inches(5.4),
            leaf_navne, leaf_values, "BAC7D9",
        )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def main():
    # --- Page config ---
    st.set_page_config(
        page_title="KU ledelseslag",
        page_icon=load_logo(),
        layout="wide",
    )
 
    col_logo, col_title = st.columns([1, 4])
    with col_logo:
        st.markdown(
            f'<img src="data:image/png;base64,{logo_base64()}" '
            f'style="max-width:180px; width:100%;">',
            unsafe_allow_html=True
        )
 
    with col_title:
        #st.title("Københavns Universitets ledelseslag")
        st.title("Personale- og lønomkostninger")

    st.markdown(
"""
Dette værktøj viser årsværk, medarbejderantal og lønomkostninger for KU's administrative enheder (
koncernenheder og campusadministrationer) og deres afdelinger. Brug menuen nedenfor til at vælge, hvilke tal figurene skal vise. 
""")
 
    # --- Data: indlæs og rul årsværk/lønomkostninger op gennem hierarkiet ---
    units = load_units()
    by_id, children_of = build_lookup_and_rollup(units)
 
    # Niveau 4 = det yderste niveau (KU=1, Enhed=2, Afdeling=3, Kontor=4)
    niveau1_navn = NIVEAUER[0]
    niveau1_ids = sorted(
        (uid for uid, u in by_id.items() if u["niveau"] == niveau1_navn),
        key=lambda uid: by_id[uid]["navn"],
    )
    
    def _nulstil_valg():
        st.session_state.pop("valgt_niveau1", None)

    with st.expander("Hvad vil du gerne se i figurene?", expanded=True):
        metric = st.radio(
            "**Vælg, hvilke tal figurene skal vise:**",
            ["Samlede lønomkostninger", "Antal medarbejdere", "Antal årsværk", "Gns. lønomkostning pr. årsværk", "Gns. lønomkostning pr. medarbejder"],
            horizontal=True,
            key="metric_valg",
        )

        vis_omraader = st.toggle(
            "Vis administrative områder",
            key="vis_omraader_toggle",
            on_change=_nulstil_valg,
        )

        omraade_valgt = None
        if vis_omraader:
            omraader = sorted(set(
                u["omraade"] for u in by_id.values()
                if u.get("niveau") == "Kontor" and u.get("omraade") is not None
            ))
            omraade_valgt = st.selectbox("**Vælg administrativt område:**", omraader, key="omraade_valg")

 
    def metric_value(uid):
        u = by_id[uid]
        if metric == "Gns. lønomkostning pr. årsværk":
            return gns_loen(by_id, uid)
        elif metric == "Gns. lønomkostning pr. medarbejder":
            return (u["lonomkostninger"] / u["medarbejdere"]) if u.get("medarbejdere") else 0
        elif metric == "Samlede lønomkostninger":
            return u["lonomkostninger"]
        elif metric == "Antal medarbejdere":
            return u["medarbejdere"]
        else:
            return u["aarsvaerk"]
 
    y_fmt = "%{y:,.1f} årsværk" if metric == "Antal årsværk" else "%{y:,.0f} kr."

    
    #col_bar1, col_bar2 = st.columns(2)
 
    # -----------------------------------------------------------------
    # Venstre: søjlediagram 1 - alle niveau 4-ledere (Kontor)
    # -----------------------------------------------------------------
    #with col_bar1:
        #st.subheader("Campusadministrationer og koncernenheder")

        #navne = [by_id[uid]["navn"] for uid in niveau1_ids]
    value_fmt = "%{x:,.1f} årsværk" if metric == "Antal årsværk" else "%{x:,.0f} kr."
    if metric == "Antal årsværk":
        value_fmt = "%{x:,.1f} årsværk"
    elif metric == "Antal medarbejdere":
        value_fmt = "%{x:,.0f} medarbejdere"
    elif metric == "Gns. lønomkostning pr. årsværk":
        value_fmt = "%{x:,.0f} kr. pr. årsværk"
    elif metric == "Gns. lønomkostning pr. medarbejder":
        value_fmt = "%{x:,.0f} kr. pr. medarbejder"
    else:
        value_fmt = "%{x:,.0f} kr."


        #if visning == "Afdelinger":
            #y = [metric_value(uid) for uid in niveau1_ids]

            #fig1 = go.Figure(go.Bar(
                #x=y,
                #y=navne,
                #orientation="h",
                #marker_color="#901A1E",
                #hovertemplate="<b>%{y}</b><br>" + value_fmt + "<extra></extra>",
            #))
            #fig1.update_layout(
                #barmode="stack",
                #margin=dict(t=40, l=10, r=10, b=50),
                #height=max(420, 28 * len(navne)),
                #xaxis_title=metric,
                #yaxis=dict(autorange="reversed"),
                #legend=dict(orientation="h", yanchor="top", y=-0.22, xanchor="left", x=0),
            #)
        #else:
            #omraade_serie, rest_serie = [], []
            #for uid in niveau1_ids:
                #o, r = _split_by_omraade(by_id, children_of, uid, omraade_valgt, metric)
                #omraade_serie.append(o)
                #rest_serie.append(r)

            #fig1 = go.Figure()
            #fig1.add_trace(go.Bar(
                #x=omraade_serie, y=navne, orientation="h", name=omraade_valgt,
                #marker=dict(color="#901A1E"),
                #hovertemplate="<b>%{y}</b><br>" + omraade_valgt + ": " + value_fmt + "<extra></extra>",
            #))
            #fig1.add_trace(go.Bar(
                #x=rest_serie, y=navne, orientation="h", name="Øvrige",
                #marker=dict(color="#E6C9CC"),
                #hovertemplate="<b>%{y}</b><br>Øvrige: " + value_fmt + "<extra></extra>",
            #))
            #fig1.update_layout(
                #barmode="stack",
                #margin=dict(t=40, l=10, r=10, b=10),
                #height=max(420, 28 * len(navne)),
                #xaxis_title=metric,
                #yaxis=dict(autorange="reversed"),
                #legend=dict(orientation="h", yanchor="top", y=-0.22, xanchor="left", x=0),
            #)

        #event = st.plotly_chart(
            #fig1,
            #key="bar_niveau1",
            #on_select="rerun",
            #selection_mode=["points"],
            #width="stretch",
        #)

        #if event and event.selection and event.selection["points"]:
            #point = event.selection["points"][0]
            #idx = point.get("point_index")
            #curve = point.get("curve_number")
            # I "Områder"-visning er trace 0 den KU-røde område-del - kun
            # klik dér skal opdatere højre diagram. I "Kontorer"-visning er
            # der kun én trace (curve altid 0), så alle klik tæller.
            #if idx is not None and curve == 0:
                #st.session_state.valgt_niveau1 = niveau1_ids[idx]
 
    # -----------------------------------------------------------------
    # Højre: søjlediagram 2 - kontorerne (yderste niveau) under den
    # valgte campusadministration/koncernenhed
    # -----------------------------------------------------------------
    #with col_bar2:
        #valgt = st.session_state.get("valgt_niveau1")
 
        #if valgt is None or valgt not in by_id:
            #st.header(" \n \n ")
            #st.header(" \n \n ")
            #st.error("Klik på en søjle til venstre for at se kontorerne under den enhed.")
        #else:
            #leaf_ids = leaves_under(children_of, valgt)
            #if visning == "Administrative områder":
                #st.subheader(f"{omraade_valgt}-andel pr. kontor under: {by_id[valgt]['navn']}")
            #else:
                #st.subheader(f"Kontorer under: {by_id[valgt]['navn']}")

            #if not leaf_ids:
                #st.info("Denne enhed har ingen underliggende kontorer i dummy-dataen.")
            #else:
                #leaf_ids = sorted(leaf_ids, key=lambda uid: by_id[uid]["navn"])
                #leaf_navne = [by_id[uid]["navn"] for uid in leaf_ids]

                #if visning == "Administrative områder":
                    #leaf_y = [
                        #metric_value(uid) if by_id[uid]["omraade"] == omraade_valgt else 0
                        #for uid in leaf_ids
                    #]
                    #leaf_farver = [
                        #"#901A1E" if by_id[uid]["omraade"] == omraade_valgt else "#E6C9CC"
                        #for uid in leaf_ids
                    #]
                #else:
                    #leaf_y = [metric_value(uid) for uid in leaf_ids]
                    #leaf_farver = "#7992b5"

                #fig2 = go.Figure(go.Bar(
                    #x=leaf_y,
                    #y=leaf_navne,
                    #orientation="h",
                    #marker_color=leaf_farver,
                    #hovertemplate="<b>%{y}</b><br>" + value_fmt + "<extra></extra>",
                #))
                #fig2.update_layout(
                    #margin=dict(t=40, l=10, r=10, b=10),
                    #height=max(420, 28 * len(leaf_navne)),
                    #xaxis_title=metric,
                    #yaxis=dict(autorange="reversed"),
                #)
                #st.plotly_chart(fig2, key="bar_kontorer", width="stretch")
    
    #st.divider()
    #st.subheader("Fuldt overblik: alle enheder og kontorer")

    overblik_niveau = st.radio(
        "**Vælg, hvilket niveau figurene skal vise:**",
        ["Niveau 3 (KE/CA)", "Niveau 4+5 (afdelinger)", "Overblik"],
        horizontal=True,
        key="overblik_niveau",
    )
    vis_enhed = overblik_niveau in ("Niveau 3 (KE/CA)", "Overblik")
    vis_kontor = overblik_niveau in ("Niveau 4+5 (afdelinger)", "Overblik")

    # Fælles x-akse-grænse på tværs af ALLE tre plots, så de er sammenlignelige.
    alle_vaerdier = []
    for uid in niveau1_ids:
        if vis_omraader:
            enhed_om, enhed_rest = _split_by_omraade(by_id, children_of, uid, omraade_valgt, metric)
            enhed_total = enhed_om + enhed_rest
        else:
            enhed_total = metric_value(uid)
        if vis_enhed:
            alle_vaerdier.append(enhed_total)
        if vis_kontor:
            for kid in children_of.get(uid, []):
                alle_vaerdier.append(metric_value(kid))
    x_maks = max(alle_vaerdier) * 1.05 if alle_vaerdier else 1

    def _unikt_navn(navn, brugte_navne):
        """
        Tilføjer et usynligt mellemrum, hvis navnet allerede optræder i
        samme plot (fx et kontor der hedder det samme som sin enhed) -
        ellers slår Plotly de to søjler sammen til én, da den kategoriske
        y-akse matcher på selve teksten, ikke listeposition.
        """
        unikt = navn
        while unikt in brugte_navne:
            unikt += " "
        brugte_navne.add(unikt)
        return unikt

    def _akse_label(metric):
        """Metricnavnet + dets enhed, til brug som x-akse-titel."""
        if metric == "Antal årsværk":
            return f"{metric} (årsværk)"
        elif metric == "Antal medarbejdere":
            return f"{metric} (medarbejdere)"
        elif metric == "Gns. lønomkostning pr. årsværk":
            return f"{metric} (kr. pr. årsværk)"
        elif metric == "Gns. lønomkostning pr. medarbejder":
            return f"{metric} (kr. pr. medarbejder)"
        return f"{metric} (kr.)"

    def _format_tal(v):
        """Tekst til visning i/ved en søjle - tomt for 0 (overskrifter, luft-rækker)."""
        if not v:
            return ""
        if metric == "Antal årsværk":
            return f"{v:,.1f}"
        elif metric == "Antal medarbejdere":
            return f"{v:,.0f}"
        return f"{v:,.0f} kr."

    if overblik_niveau == "Niveau 3 (KE/CA)":
        # Ét samlet plot. Klik på en enheds-søjle folder dens kontorer ud lige
        # under den - flere enheder kan være udfoldet samtidig.
        if "niveau3_udvidet" not in st.session_state:
            st.session_state.niveau3_udvidet = set()

        if vis_omraader:
            enheder_at_vise = []
            for uid in niveau1_ids:
                om_v, rest_v = _split_by_omraade(by_id, children_of, uid, omraade_valgt, metric)
                enheder_at_vise.append((uid, om_v, rest_v))
        else:
            enheder_at_vise = [(uid, metric_value(uid), 0) for uid in niveau1_ids]

        navne, om_vaerdier, rest_vaerdier, om_kleur, klik_uid = [], [], [], [], []
        for uid, om_v, rest_v in enheder_at_vise:
            navne.append(by_id[uid]["navn"])
            om_vaerdier.append(om_v)
            rest_vaerdier.append(rest_v)
            om_kleur.append("#901A1E")
            klik_uid.append(uid)

            if uid in st.session_state.niveau3_udvidet:
                if vis_omraader:
                    kontor_ids = [
                        kid for kid in children_of.get(uid, [])
                        if by_id[kid]["omraade"] == omraade_valgt
                    ]
                else:
                    kontor_ids = children_of.get(uid, [])
                kontor_ids = sorted(kontor_ids, key=metric_value, reverse=True)
                for kid in kontor_ids:
                    navne.append(by_id[kid]["navn"])
                    om_vaerdier.append(metric_value(kid))
                    rest_vaerdier.append(0)
                    om_kleur.append("#7992b5")
                    klik_uid.append(None)

        brugte_navne_n3 = set()
        navne, om_vaerdier, rest_vaerdier, om_kleur, klik_uid = [], [], [], [], []

        for uid, om_v, rest_v in enheder_at_vise:
            navne.append(by_id[uid]["navn"])
            om_vaerdier.append(om_v)
            rest_vaerdier.append(rest_v)
            om_kleur.append("#901A1E")
            klik_uid.append(uid)

            if uid in st.session_state.niveau3_udvidet:
                if vis_omraader:
                    alle_kontorer = children_of.get(uid, [])
                    kontor_ids_match = sorted(
                        (kid for kid in alle_kontorer if by_id[kid]["omraade"] == omraade_valgt),
                        key=metric_value, reverse=True,
                    )
                    kontor_ids_oevrige = sorted(
                        (kid for kid in alle_kontorer if by_id[kid]["omraade"] != omraade_valgt),
                        key=metric_value, reverse=True,
                    )
                    for kid in kontor_ids_match:
                        navne.append(by_id[kid]["navn"])
                        om_vaerdier.append(metric_value(kid))
                        rest_vaerdier.append(0)
                        om_kleur.append("#7992b5")
                        klik_uid.append(None)
                    for kid in kontor_ids_oevrige:
                        navne.append(by_id[kid]["navn"])
                        om_vaerdier.append(metric_value(kid))
                        rest_vaerdier.append(0)
                        om_kleur.append("#cad4e2")
                        klik_uid.append(None)
                else:
                    kontor_ids = sorted(children_of.get(uid, []), key=metric_value, reverse=True)
                    for kid in kontor_ids:
                        navne.append(by_id[kid]["navn"])
                        om_vaerdier.append(metric_value(kid))
                        rest_vaerdier.append(0)
                        om_kleur.append("#7992b5")
                        klik_uid.append(None)

        brugte_navne_n3 = set()
        navne = [_unikt_navn(n, brugte_navne_n3) for n in navne]

        fig_niveau3 = go.Figure()
        fig_niveau3.add_trace(go.Bar(
            x=om_vaerdier,
            y=navne,
            orientation="h",
            name=omraade_valgt if vis_omraader else metric,
            marker_color=om_kleur,
            marker_line_color="white",
            marker_line_width=1,
            text=[_format_tal(v) for v in om_vaerdier],
            textposition="auto",
            hovertemplate="<b>%{y}</b><br>" + value_fmt + "<extra></extra>",
        ))
        if vis_omraader:
            fig_niveau3.add_trace(go.Bar(
                x=rest_vaerdier,
                y=navne,
                orientation="h",
                name="Øvrige",
                marker_color="#E6C9CC",
                marker_line_color="white",
                marker_line_width=1,
                text=[_format_tal(v) for v in rest_vaerdier],
                textposition="auto",
                hovertemplate="<b>%{y}</b><br>Øvrige: " + value_fmt + "<extra></extra>",
            ))
        fig_niveau3.update_layout(
            barmode="stack",
            title=f"{metric} for niveau 3",
            margin=dict(t=60, l=10, r=10, b=10),
            height=max(420, 20 * len(navne)),
            xaxis=dict(title=_akse_label(metric)),
            yaxis=dict(autorange="reversed"),
            bargap=0,
        )

        event_n3 = st.plotly_chart(
            fig_niveau3,
            key="overblik_niveau3_samlet",
            on_select="rerun",
            selection_mode=["points"],
            width="stretch",
        )

        if event_n3 and event_n3.selection and event_n3.selection["points"]:
            idx = event_n3.selection["points"][0].get("point_index")
            if idx is not None and idx < len(klik_uid) and klik_uid[idx] is not None:
                klikket_uid = klik_uid[idx]
                if klikket_uid in st.session_state.niveau3_udvidet:
                    st.session_state.niveau3_udvidet.discard(klikket_uid)
                else:
                    st.session_state.niveau3_udvidet.add(klikket_uid)
                st.rerun()
        
        st.caption("Klik på en søjle ovenfor for at folde dens afdelinger ud.")

    elif overblik_niveau == "Niveau 4+5 (afdelinger)":
        # Samme klyngestruktur som Begge niveauer (CA/KE øverst i hver
        # gruppe), men uden en rigtig enheds-søjle - kun en tom
        # "overskrift"-søjle med enhedens navn som label, efterfulgt af
        # dens kontorer.
        CA_RAEKKEFOELGE = [
            "Campusadministration Frederiksberg+",
            "Campusadministration Nørre",
            "Campusadministration Søndre",
        ]
        ca_ids = sorted(
            (uid for uid in niveau1_ids if uid in CA_RAEKKEFOELGE),
            key=lambda uid: CA_RAEKKEFOELGE.index(uid),
        )
        ovrige_ids = [uid for uid in niveau1_ids if uid not in CA_RAEKKEFOELGE]

        chunk_n4 = -(-len(ovrige_ids) // 3)  # oprund
        ovrige_grupper_n4 = [ovrige_ids[i:i + chunk_n4] for i in range(0, len(ovrige_ids), chunk_n4)]
        while len(ovrige_grupper_n4) < 3:
            ovrige_grupper_n4.append([])
        while len(ca_ids) < 3:
            ca_ids.append(None)

        grupper_n4 = [
            ([ca] if ca is not None else []) + ovrige
            for ca, ovrige in zip(ca_ids, ovrige_grupper_n4)
        ]

        kolonner_n4 = st.columns(3)
        for g_idx, gruppe in enumerate(grupper_n4):
            navne, vaerdier, om_kleur, overskrift_navne = [], [], [], []
            brugte_navne_n4 = set()

            for uid in gruppe:
                if vis_omraader:
                    alle_kontorer = children_of.get(uid, [])
                    kontor_ids_match = sorted(
                        (kid for kid in alle_kontorer if by_id[kid]["omraade"] == omraade_valgt),
                        key=metric_value, reverse=True,
                    )
                    kontor_ids_oevrige = sorted(
                        (kid for kid in alle_kontorer if by_id[kid]["omraade"] != omraade_valgt),
                        key=metric_value, reverse=True,
                    )
                    kontor_ids_samlet = kontor_ids_match + kontor_ids_oevrige
                    farve_pr_kontor = ["#7992b5"] * len(kontor_ids_match) + ["#DCE3EC"] * len(kontor_ids_oevrige)
                else:
                    kontor_ids_samlet = sorted(children_of.get(uid, []), key=metric_value, reverse=True)
                    farve_pr_kontor = ["#7992b5"] * len(kontor_ids_samlet)

                if not kontor_ids_samlet:
                    continue  # ingen kontorer at vise for denne enhed

                if navne:  # luft foer alle overskrifter undtagen den foerste i plottet
                    navne.append(_unikt_navn(" ", brugte_navne_n4))
                    vaerdier.append(0)
                    om_kleur.append("rgba(0,0,0,0)")

                # Overskrift: tom søjle - navnet vises som annotation INDE i
                # plottet i stedet for som akse-label, se tickvals nedenfor.
                overskrift_id = _unikt_navn(by_id[uid]["navn"], brugte_navne_n4)
                navne.append(overskrift_id)
                vaerdier.append(0)
                om_kleur.append("rgba(0,0,0,0)")
                overskrift_navne.append((overskrift_id, by_id[uid]["navn"]))

                for kid, farve in zip(kontor_ids_samlet, farve_pr_kontor):
                    navne.append(_unikt_navn(by_id[kid]["navn"], brugte_navne_n4))
                    vaerdier.append(metric_value(kid))
                    om_kleur.append(farve)

            if not navne:
                continue

            fig_niveau4 = go.Figure(go.Bar(
                x=vaerdier,
                y=navne,
                orientation="h",
                marker_color=om_kleur,
                marker_line_color="white",
                marker_line_width=1,
                text=[_format_tal(v) for v in vaerdier],
                textposition="auto",
                hovertemplate="<b>%{y}</b><br>" + value_fmt + "<extra></extra>",
            ))

            # Kun kontor-rækkerne skal have en akse-label i margenen -
            # overskrifterne vises i stedet som annotationer inde i plottet.
            overskrift_ids = {oid for oid, _ in overskrift_navne}
            tick_rækker = [n for n in navne if n not in overskrift_ids]

            for overskrift_id, visningsnavn in overskrift_navne:
                fig_niveau4.add_annotation(
                    x=-0.0, y=overskrift_id,
                    xref="paper", yref="y",
                    text=visningsnavn,
                    showarrow=False,
                    xanchor="center",
                    font=dict(size=12, color="#838697"),  # samme font som akse-labels
                )

            fig_niveau4.update_layout(
                title=f"{metric} for niveau 4+5" if g_idx == 0 else " ",
                margin=dict(t=60, l=20, r=10, b=10),
                height=max(160, 20 * len(navne) + 60),
                xaxis=dict(title=_akse_label(metric)),
                yaxis=dict(autorange="reversed", tickmode="array", tickvals=tick_rækker, ticktext=tick_rækker),
                bargap=0,
            )
            with kolonner_n4[g_idx]:
                st.plotly_chart(fig_niveau4, key=f"overblik_niveau4_plot_{g_idx}", width="stretch")
    else:

        # Del de 13 enheder i tre nogenlunde lige store, sammenhængende grupper -
        # én gruppe pr. kolonne/plot.
        CA_RAEKKEFOELGE = [
            "Campusadministration Frederiksberg+",
            "Campusadministration Nørre",
            "Campusadministration Søndre",
        ]
        ca_ids = sorted(
            (uid for uid in niveau1_ids if uid in CA_RAEKKEFOELGE),
            key=lambda uid: CA_RAEKKEFOELGE.index(uid),
        )
        ovrige_ids = [uid for uid in niveau1_ids if uid not in CA_RAEKKEFOELGE]

        chunk = -(-len(ovrige_ids) // 3)  # oprund
        ovrige_grupper = [ovrige_ids[i:i + chunk] for i in range(0, len(ovrige_ids), chunk)]
        while len(ovrige_grupper) < 3:
            ovrige_grupper.append([])
        while len(ca_ids) < 3:
            ca_ids.append(None)

        grupper = [
            ([ca] if ca is not None else []) + ovrige
            for ca, ovrige in zip(ca_ids, ovrige_grupper)
        ]

        kolonner = st.columns(3)
        for g_idx, gruppe in enumerate(grupper):
            navne, om_vaerdier, rest_vaerdier, om_kleur = [], [], [], []
            brugte_navne = set()

            for uid in gruppe:
                if vis_omraader:
                    alle_kontorer = children_of.get(uid, [])
                    kontor_ids_match = sorted(
                        (kid for kid in alle_kontorer if by_id[kid]["omraade"] == omraade_valgt),
                        key=metric_value, reverse=True,
                    )
                    kontor_ids_oevrige = sorted(
                        (kid for kid in alle_kontorer if by_id[kid]["omraade"] != omraade_valgt),
                        key=metric_value, reverse=True,
                    )
                    kontor_ids = kontor_ids_match + kontor_ids_oevrige
                    farve_pr_kontor = ["#7992b5"] * len(kontor_ids_match) + ["#DCE3EC"] * len(kontor_ids_oevrige)
                    if not kontor_ids:
                        continue  # ingen kontorer overhovedet under denne enhed
                    enhed_om, enhed_rest = _split_by_omraade(by_id, children_of, uid, omraade_valgt, metric)
                else:
                    enhed_om, enhed_rest = metric_value(uid), 0
                    kontor_ids = sorted(children_of.get(uid, []), key=metric_value, reverse=True)
                    farve_pr_kontor = ["#7992b5"] * len(kontor_ids)

                if vis_enhed:
                    if navne and vis_kontor:
                        navne.append(_unikt_navn(" ", brugte_navne))
                        om_vaerdier.append(0)
                        rest_vaerdier.append(0)
                        om_kleur.append("rgba(0,0,0,0)")
                    navne.append(_unikt_navn(by_id[uid]["navn"], brugte_navne))
                    om_vaerdier.append(enhed_om)
                    rest_vaerdier.append(enhed_rest)
                    om_kleur.append("#901A1E")

                if vis_kontor:
                    for kid, farve in zip(kontor_ids, farve_pr_kontor):
                        navne.append(_unikt_navn(by_id[kid]["navn"], brugte_navne))
                        om_vaerdier.append(metric_value(kid))
                        rest_vaerdier.append(0)
                        om_kleur.append(farve)

            if not navne:
                continue

            fig_overblik = go.Figure()
            fig_overblik.add_trace(go.Bar(
                x=om_vaerdier,
                y=navne,
                orientation="h",
                name=omraade_valgt if vis_omraader else metric,
                marker_color=om_kleur,
                marker_line_color="white",
                marker_line_width=1,
                text=[_format_tal(v) for v in om_vaerdier],
                textposition="auto",
                hovertemplate="<b>%{y}</b><br>" + value_fmt + "<extra></extra>",
            ))
            if vis_omraader:
                fig_overblik.add_trace(go.Bar(
                    x=rest_vaerdier,
                    y=navne,
                    orientation="h",
                    name="Øvrige",
                    marker_color="#E6C9CC",
                    marker_line_color="white",
                    marker_line_width=1,
                    text=[_format_tal(v) for v in rest_vaerdier],
                    textposition="auto",
                    hovertemplate="<b>%{y}</b><br>Øvrige: " + value_fmt + "<extra></extra>",
                ))
            fig_overblik.update_layout(
                barmode="stack",
                title=f"{metric} for alle tre niveauer" if g_idx == 0 else "",
                margin=dict(t=60, l=10, r=10, b=70 if g_idx == 1 else 10),
                height=max(160, 20 * len(navne) + 60),
                xaxis=dict(range=[0, x_maks], title=_akse_label(metric)),
                yaxis=dict(autorange="reversed"),
                bargap=0,
                showlegend=(g_idx == 4),
                legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5),
            )

            with kolonner[g_idx]:
                st.plotly_chart(fig_overblik, key=f"overblik_plot_{g_idx}", width="stretch")
    #st.divider()

    #st.subheader("Se én KE/CA i detaljer")
    #st.markdown(
#"""
#Vælg af listen nedenfor, hvilken campusadministration eller koncernenhed du vil se nærmere på. 
#""")

    #zoom_navne = [by_id[uid]["navn"] for uid in niveau1_ids]
    #zoom_valgt_navn = st.selectbox("**Vælg enhed:**", zoom_navne, key="zoom_valg")
    #zoom_valgt_uid = next(uid for uid in niveau1_ids if by_id[uid]["navn"] == zoom_valgt_navn)

    #leaf_ids = leaves_under(children_of, zoom_valgt_uid)
    #if vis_omraader:
        #titeltekst = f"{omraade_valgt}-andel pr. kontor under: {by_id[zoom_valgt_uid]['navn']}"
    #else:
        #titeltekst = f"Kontorer under: {by_id[zoom_valgt_uid]['navn']}"

    #if not leaf_ids:
        #st.info("Denne enhed har ingen underliggende kontorer.")
    #else:
        #leaf_ids = sorted(leaf_ids, key=metric_value, reverse=True)
        #leaf_navne = [by_id[uid]["navn"] for uid in leaf_ids]

        #if vis_omraader:
            #leaf_y = [
                #metric_value(uid) if by_id[uid]["omraade"] == omraade_valgt else 0
                #for uid in leaf_ids
            #]
            #leaf_farver = [
                #"#901A1E" if by_id[uid]["omraade"] == omraade_valgt else "#E6C9CC"
                #for uid in leaf_ids
            #]
        #else:
            #leaf_y = [metric_value(uid) for uid in leaf_ids]
            #leaf_farver = "#7992b5"

        #fig_zoom = go.Figure(go.Bar(
            #x=leaf_y,
            #y=leaf_navne,
            #orientation="h",
            #marker_color=leaf_farver,
            #marker_line_color="white",
            #marker_line_width=1,
            #hovertemplate="<b>%{y}</b><br>" + value_fmt + "<extra></extra>",
        #))
        #fig_zoom.update_layout(
            #margin=dict(t=60, l=10, r=10, b=10),
            #height=max(420, 28 * len(leaf_navne)),
            #title=titeltekst,
            #xaxis_title=_akse_label(metric),
            #yaxis=dict(autorange="reversed"),
        #)
        #st.plotly_chart(fig_zoom, key="bar_zoom", width="stretch")



    st.markdown(" \n ")

    if st.button("Generér PowerPoint med alle KE/CA"):
        with st.spinner("Bygger PowerPoint..."):
            pptx_buf = build_full_pptx(by_id, children_of, niveau1_ids, metric, metric_value)
        st.download_button(
            "Download PowerPoint",
            data=pptx_buf,
            file_name="ledelseslag_alle_enheder.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    
    # Footer
    st.markdown(f"""
<hr style="margin-top: 50px;">
<div style="text-align:center; color:#666; font-size: 0.9em;">
  REKSTAB Analyse · Amanda Schramm Petersen · <a href="mailto:ascp@adm.ku.dk">ascp@adm.ku.dk</a>
  · opdateret 1. september 2026
</div>
""", unsafe_allow_html=True)
 
if __name__ == "__main__":
    main()